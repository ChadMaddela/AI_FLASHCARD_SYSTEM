import io
import logging
import os

import docx
import PyPDF2
import pytesseract
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

IMAGE_EXTENSIONS = {".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png"}


def extract_text_from_file(uploaded_file):
    """Extract text depending on file type from an uploaded file object or path."""
    ext = os.path.splitext(uploaded_file.name)[1].lower()

    if ext == ".pdf":
        text = ""
        reader = PyPDF2.PdfReader(uploaded_file)
        for page in reader.pages:
            text += page.extract_text() or ""
        return text

    elif ext == ".docx":
        file_stream = io.BytesIO(uploaded_file.read())
        doc = docx.Document(file_stream)
        return "\n".join([para.text for para in doc.paragraphs])

    elif ext == ".pptx":
        file_stream = io.BytesIO(uploaded_file.read())
        prs = Presentation(file_stream)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    elif ext in IMAGE_EXTENSIONS:
        file_stream = io.BytesIO(uploaded_file.read())
        img = Image.open(file_stream)
        return pytesseract.image_to_string(img)

    else:
        raise ValueError(f"Unsupported file type: {ext}")


def extract_images_from_bytes(file_bytes: bytes, filename: str, max_images: int = 12):
    """
    Returns a list of (mime_type, image_bytes) tuples extracted from the file, capped at
    max_images. Never raises — any extraction failure just yields an empty list, so a material
    with problematic embedded images simply falls back to text-only flashcard generation.
    """
    ext = os.path.splitext(filename)[1].lower()
    images = []

    try:
        if ext == ".pdf":
            reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            for page in reader.pages:
                for img in getattr(page, "images", []):
                    images.append(("image/png", img.data))
                    if len(images) >= max_images:
                        return images

        elif ext == ".docx":
            doc = docx.Document(io.BytesIO(file_bytes))
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    images.append((rel.target_part.content_type, rel.target_part.blob))
                    if len(images) >= max_images:
                        return images

        elif ext == ".pptx":
            prs = Presentation(io.BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        images.append((shape.image.content_type, shape.image.blob))
                        if len(images) >= max_images:
                            return images

        elif ext in IMAGE_EXTENSIONS:
            images.append((IMAGE_EXTENSIONS[ext], file_bytes))

    except Exception:
        logger.exception("Image extraction failed for %s; continuing without images.", filename)
        return []

    return images


def downscale_image(mime_type: str, data: bytes, max_dimension: int = 1024) -> tuple:
    """Shrinks+re-encodes an image for a cheaper/faster Gemini request. Never used for storage/display."""
    try:
        img = Image.open(io.BytesIO(data))
        img.thumbnail((max_dimension, max_dimension))
        buf = io.BytesIO()
        img.convert("RGB").save(buf, format="JPEG", quality=80)
        return "image/jpeg", buf.getvalue()
    except Exception:
        logger.warning("Could not downscale an extracted image; sending original bytes to Gemini instead.")
        return mime_type, data
