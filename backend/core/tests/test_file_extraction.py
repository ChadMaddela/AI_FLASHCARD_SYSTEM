import os
import django
import pytest
from io import BytesIO
from PIL import Image

os.environ.setdefault('DJANGO_SETTINGS_MODULE', 'config.settings')
django.setup()

from core.file_extraction import extract_images_from_bytes, downscale_image


def _make_png_bytes(size=(10, 10), color=(255, 0, 0)):
    buf = BytesIO()
    Image.new("RGB", size, color=color).save(buf, format="PNG")
    return buf.getvalue()


def _build_docx_with_image(n=1):
    from docx import Document
    doc = Document()
    doc.add_paragraph("Some source text.")
    for _ in range(n):
        img_buf = BytesIO(_make_png_bytes())
        doc.add_picture(img_buf)
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _build_pptx_with_image(n=1):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for _ in range(n):
        img_buf = BytesIO(_make_png_bytes())
        slide.shapes.add_picture(img_buf, Inches(1), Inches(1))
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_extract_images_from_docx():
    data = _build_docx_with_image()
    images = extract_images_from_bytes(data, "material.docx")
    assert len(images) == 1
    mime, img_bytes = images[0]
    assert mime.startswith("image/")
    assert len(img_bytes) > 0


def test_extract_images_from_pptx():
    data = _build_pptx_with_image()
    images = extract_images_from_bytes(data, "material.pptx")
    assert len(images) == 1
    mime, img_bytes = images[0]
    assert "image" in mime
    assert len(img_bytes) > 0


def test_extract_images_passthrough_for_plain_image_upload():
    data = _make_png_bytes()
    images = extract_images_from_bytes(data, "photo.png")
    assert images == [("image/png", data)]


def test_extract_images_returns_empty_for_unsupported_extension():
    assert extract_images_from_bytes(b"whatever", "notes.txt") == []


def test_extract_images_respects_max_images_cap():
    data = _build_pptx_with_image(n=5)
    images = extract_images_from_bytes(data, "material.pptx", max_images=2)
    assert len(images) == 2


def test_extract_images_never_raises_on_corrupt_file():
    # A .docx extension but garbage bytes — must degrade to [] instead of propagating an exception,
    # so a single problematic material never blocks the rest of flashcard generation.
    assert extract_images_from_bytes(b"not a real docx file", "broken.docx") == []


def test_downscale_image_produces_valid_jpeg_under_cap():
    data = _make_png_bytes(size=(2000, 2000))
    mime, downscaled = downscale_image("image/png", data, max_dimension=500)
    assert mime == "image/jpeg"
    result_img = Image.open(BytesIO(downscaled))
    assert max(result_img.size) <= 500
    assert result_img.format == "JPEG"


def test_downscale_image_falls_back_to_original_on_invalid_data():
    mime, data = downscale_image("image/png", b"not an image", max_dimension=500)
    assert data == b"not an image"
    assert mime == "image/png"
