import json
import logging
import os
import io
from django.shortcuts import get_object_or_404
from django.db import transaction
from django.core.files.storage import default_storage
from rest_framework.decorators import api_view, permission_classes, parser_classes
from rest_framework.permissions import IsAuthenticated
from rest_framework.response import Response
from rest_framework import status
from rest_framework.parsers import MultiPartParser, FormParser, JSONParser
from google import genai
from google.genai import types
from config import settings
from .models import Material, Flashcard, StudentPerformance, User
from .spacing_engine import process_student_response, get_next_cards_for_student
from .serializers import UserSerializer, FlashcardSerializer, MaterialSerializer, MaterialCreateSerializer
import docx
import PyPDF2
from pptx import Presentation
from PIL import Image
import pytesseract

# Explicit import of Supabase client client wrapper layout mapping variables
from supabase import create_client  

logger = logging.getLogger(__name__)
client = genai.Client(api_key=settings.GEMINI_API_KEY)

# Initialize Supabase Connection
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")
supabase = create_client(SUPABASE_URL, SUPABASE_KEY)


def extract_text_from_file(uploaded_file):
    """Extract text depending on file type from an uploaded file object or path."""
    ext = os.path.splitext(uploaded_file.name)[1].lower()

    if ext == ".pdf":
        text = ""
        reader = PyPDF2.PdfReader(uploaded_file)
        for page in reader.pages:
            text += page.extract_text() or ""
        return text

    elif ext == ".docx":
        file_stream = io.BytesIO(uploaded_file.read())
        doc = docx.Document(file_stream)
        return "\n".join([para.text for para in doc.paragraphs])

    elif ext == ".pptx":
        file_stream = io.BytesIO(uploaded_file.read())
        prs = Presentation(file_stream)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    elif ext in [".jpg", ".jpeg", ".png"]:
        file_stream = io.BytesIO(uploaded_file.read())
        img = Image.open(file_stream)
        return pytesseract.image_to_string(img)

    else:
        raise ValueError(f"Unsupported file type: {ext}")


@api_view(["GET"])
@permission_classes([IsAuthenticated])
def student_card_queue(request, material_id):
    """Retrieve up to 10 cards for review based on mastery tracking and an optional sub-topic filter."""
    material = get_object_or_404(Material, id=material_id)
    student = request.user
    
    # 🎛️ Extract sub_topic filter from request query string
    sub_topic_filter = request.query_params.get("sub_topic", None)

    # Pull basic list using your spacing configuration tracking function
    selected_cards = get_next_cards_for_student(student, material=material, limit=10)

    # If queue is empty or a topic filter is actively selected, construct the filtered dataset safely
    if not selected_cards or len(selected_cards) == 0:
        card_queryset = Flashcard.objects.filter(material=material)
        if sub_topic_filter:
            card_queryset = card_queryset.filter(sub_topic=sub_topic_filter)
        selected_cards = card_queryset.order_by("id")[:10]
    elif sub_topic_filter:
        # Re-fetch targeted sub_topics to ensure we load a full set of 10 items for this specific topic
        card_queryset = Flashcard.objects.filter(material=material, sub_topic=sub_topic_filter)
        selected_cards = card_queryset.order_by("id")[:10]

    payload = []
    for c in selected_cards:
        perf = StudentPerformance.objects.filter(student=student, flashcard=c).first()
        payload.append({
            "id": c.id,
            "question": c.question,
            "choices": {
                "A": c.choice_a,
                "B": c.choice_b,
                "C": c.choice_c,
                "D": c.choice_d
            },
            "sub_topic": c.sub_topic,
            "current_mastery_level": perf.mastery_level if perf else 0
        })

    # 🗂️ Pull unique string tags directly from the material to populate navigation pills dynamically
    all_sub_topics = list(
        Flashcard.objects.filter(material=material)
        .values_list("sub_topic", flat=True)
        .distinct()
    )
    # Clean up empty strings or null entries safely
    all_sub_topics = [topic for topic in all_sub_topics if topic]

    return Response({
        "queue": payload,
        "available_topics": all_sub_topics
    }, status=status.HTTP_200_OK)

@api_view(["POST"])
@permission_classes([IsAuthenticated])
@parser_classes([MultiPartParser, FormParser, JSONParser])
def teacher_upload_material(request):
    """Teachers upload lecture materials (text or file) → Gemini generates flashcards."""
    if request.user.role != "TEACHER":
        return Response({"error": "Unauthorized. Only teachers can upload materials."}, 
                        status=status.HTTP_403_FORBIDDEN)

    title = request.data.get("title")
    content_text = request.data.get("content")
    uploaded_file = request.FILES.get("file")

    if not title and not uploaded_file:
        return Response({"error": "Provide either 'title' + 'content' or upload a file."},
                        status=status.HTTP_400_BAD_REQUEST)

    try:
        if uploaded_file:
            content_text = extract_text_from_file(uploaded_file)

        if not content_text or not content_text.strip():
            return Response({"error": "No text could be extracted from the file/content input."},
                            status=status.HTTP_400_BAD_REQUEST)

        # 1. Save base metadata record to database
        material = Material.objects.create(
            title=title or uploaded_file.name,
            description=f"Uploaded material processed. Size: {len(content_text)} chars.",
            content_text=content_text,
            uploaded_by=request.user
        )

        if uploaded_file:
            uploaded_file.seek(0)
            raw_file_bytes = uploaded_file.read()
            
            file_ext = os.path.splitext(uploaded_file.name)[1].lower()
            clean_filename = f"material_{material.id}{file_ext}"
            supabase_storage_path = f"materials/{clean_filename}"
            
            bucket_name = "materials"
            
            # 2. Upload file to storage (Will bypass RLS cleanly using the Service Role Key!)
            supabase.storage.from_(bucket_name).upload(
                path=supabase_storage_path,
                file=raw_file_bytes,
                file_options={"content-type": getattr(uploaded_file, "content_type", "application/octet-stream")}
            )
            
            # Fetch the public URL response object
            public_url_response = supabase.storage.from_(bucket_name).get_public_url(supabase_storage_path)
            
            # ✅ SAFE EXTRACTION: Safely convert response variants into a pure text URL string
            if isinstance(public_url_response, dict):
                generated_url = public_url_response.get("publicUrl", "")
            elif hasattr(public_url_response, "public_url"):
                generated_url = public_url_response.public_url
            else:
                generated_url = str(public_url_response)

            # Assign text value and save cleanly
            material.file_url = generated_url
            material.save(update_fields=["file_url"])

        # 3. Fire Gemini LLM call
        prompt = (
            f"Generate 5–10 multiple-choice questions (MCQs) from the following material. "
            f"Each question must have 4 choices (A–D), identify the correct choice, and state the sub-topic.\n\n"
            f"Study Material:\n{content_text}"
        )

        response_schema = types.Schema(
            type=types.Type.ARRAY,
            items=types.Schema(
                type=types.Type.OBJECT,
                properties={
                    "question": types.Schema(type=types.Type.STRING),
                    "choice_a": types.Schema(type=types.Type.STRING),
                    "choice_b": types.Schema(type=types.Type.STRING),
                    "choice_c": types.Schema(type=types.Type.STRING),
                    "choice_d": types.Schema(type=types.Type.STRING),
                    "correct_choice": types.Schema(type=types.Type.STRING),
                    "sub_topic": types.Schema(type=types.Type.STRING),
                },
                required=["question", "choice_a", "choice_b", "choice_c", "choice_d", "correct_choice", "sub_topic"]
            )
        )

        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt,
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=response_schema,
                temperature=0.2
            )
        )

        raw_output = response.text
        if not raw_output:
            raise ValueError("Gemini AI generated an empty response payload.")

        generated_cards = json.loads(raw_output)

        flashcards_to_create = [
            Flashcard(
                material=material,
                question=item["question"],
                choice_a=item["choice_a"],
                choice_b=item["choice_b"],
                choice_c=item["choice_c"],
                choice_d=item["choice_d"],
                correct_choice=item["correct_choice"].upper().strip(),
                sub_topic=item["sub_topic"]
            )
            for item in generated_cards
        ]

        with transaction.atomic():
            Flashcard.objects.bulk_create(flashcards_to_create)

        return Response({
            "message": "Material uploaded and flashcards created via Gemini AI.",
            "material_id": material.id,
            "file_url": material.file_url,
            "flashcards_count": len(flashcards_to_create)
        }, status=status.HTTP_201_CREATED)

    except Exception as e:
        logger.exception("Failed to process upload")
        return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(["POST"])
@permission_classes([IsAuthenticated])
def submit_answer(request):
    """Process student answer → update mastery level + analytics."""
    flashcard_id = request.data.get("flashcard_id")
    selected_choice = request.data.get("selected_choice")

    if not flashcard_id or not selected_choice:
        return Response({"error": "Fields 'flashcard_id' and 'selected_choice' are required."},
                        status=status.HTTP_400_BAD_REQUEST)

    flashcard = get_object_or_404(Flashcard, id=flashcard_id)
    student = request.user

    perf = process_student_response(student, flashcard, selected_choice)
    is_correct = (selected_choice.strip().upper() == flashcard.correct_choice.upper())
    accuracy_rate = perf.accuracy_percentage

    return Response({
        "is_correct": is_correct,
        "correct_answer": flashcard.correct_choice,
        "analytics": {
            "updated_mastery_level": perf.mastery_level,
            "total_attempts": perf.attempts_count,
            "accuracy_percentage": accuracy_rate,
            "sub_topic_tracked": flashcard.sub_topic
        }
    }, status=status.HTTP_200_OK)


@api_view(["GET"])
@permission_classes([IsAuthenticated])
def user_me(request):
    """Return full info about the currently authenticated user, including their role."""
    serializer = UserSerializer(request.user)
    return Response(serializer.data)


@api_view(["GET"])
@permission_classes([IsAuthenticated])
def list_flashcards(request):
    """List all flashcards (teacher view)."""
    if request.user.role != "TEACHER":
        return Response({"error": "Access denied."}, status=status.HTTP_403_FORBIDDEN)
        
    flashcards = Flashcard.objects.all().order_by("-id")
    serializer = FlashcardSerializer(flashcards, many=True)
    return Response({"flashcards": serializer.data}, status=status.HTTP_200_OK)


@api_view(["PUT"])
@permission_classes([IsAuthenticated])
def update_flashcard(request, flashcard_id):
    """Update a flashcard (teacher edit)."""
    if request.user.role != "TEACHER":
        return Response({"error": "Access denied."}, status=status.HTTP_403_FORBIDDEN)

    flashcard = get_object_or_404(Flashcard, id=flashcard_id)
    serializer = FlashcardSerializer(flashcard, data=request.data, partial=True)
    if serializer.is_valid():
        serializer.save()
        return Response(serializer.data, status=status.HTTP_200_OK)
    return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)


@api_view(["DELETE"])
@permission_classes([IsAuthenticated])
def delete_flashcard(request, flashcard_id):
    """Delete a flashcard (teacher delete)."""
    if request.user.role != "TEACHER":
        return Response({"error": "Access denied."}, status=status.HTTP_403_FORBIDDEN)

    flashcard = get_object_or_404(Flashcard, id=flashcard_id)
    flashcard.delete()
    return Response({"message": "Flashcard deleted."}, status=status.HTTP_204_NO_CONTENT)


@api_view(["GET"])
@permission_classes([IsAuthenticated])
def list_materials(request):
    """List lecture materials cleanly using MaterialSerializer to preserve file_url."""
    user = request.user

    if user.role == "TEACHER":
        materials = Material.objects.filter(uploaded_by=user).order_by("-created_at")
    elif user.role == "STUDENT":
        materials = Material.objects.all().order_by("-created_at")
    else:
        materials = Material.objects.none()

    serializer = MaterialSerializer(materials, many=True)
    return Response(serializer.data, status=status.HTTP_200_OK)


@api_view(["DELETE"])
@permission_classes([IsAuthenticated])
def delete_material(request, material_id):
    """Delete lecture material record, cascaded flashcards, and remove asset from Supabase bucket."""
    if request.user.role != "TEACHER":
        return Response({"error": "Unauthorized. Only teachers can manage storage lifecycles."}, 
                        status=status.HTTP_403_FORBIDDEN)

    material = get_object_or_404(Material, id=material_id)

    try:
        # If a file link exists, parse out its name to purge the asset from Supabase storage
        if material.file_url:
            bucket_name = "materials"
            # Target the trailing segment after the bucket key identification route string
            filename = material.file_url.split(f"storage/v1/object/public/{bucket_name}/")[-1]
            
            # If splitting didn't target cleanly, parse standard fallback matching rule layout
            if "materials/" not in filename:
                filename = f"materials/{filename.split('/')[-1]}"
                
            try:
                supabase.storage.from_(bucket_name).remove([filename])
            except Exception as bucket_err:
                logger.warning(f"Storage bucket drop warning or bypassed file skip trace: {str(bucket_err)}")

        # Drop DB object row cleanly (Cascades down to clean performance trackers/flashcards automatically)
        material.delete()
        return Response({"message": "Material and storage bucket payload cleanly removed."}, status=status.HTTP_200_OK)

    except Exception as e:
        logger.exception("Failed to execute material storage destruction lifecycles")
        return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)


@api_view(['GET'])
@permission_classes([IsAuthenticated])
def material_flashcards(request, material_id):
    flashcards = Flashcard.objects.filter(material_id=material_id)
    serializer = FlashcardSerializer(flashcards, many=True)
    return Response(serializer.data)


@api_view(['POST'])
@permission_classes([IsAuthenticated])
def create_flashcard(request, material_id):
    serializer = FlashcardSerializer(data={**request.data, "material": material_id})
    serializer.is_valid(raise_exception=True)
    serializer.save()
    return Response(serializer.data, status=201)